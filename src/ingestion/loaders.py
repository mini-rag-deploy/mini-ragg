# src/ingestion/loaders.py
"""
Multi-format document loader.
Supports: PDF, Word (.docx), PowerPoint (.pptx), Text (.txt), Images, URLs.

Handles:
- Corrupted / password-protected files (graceful skip)
- Encoding issues in text extraction
- Empty / unreadable pages
- Mixed-language content (Arabic + English)
- Large files (streaming page-by-page)
"""

from __future__ import annotations

import io
import logging
import re
import unicodedata
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterator, List, Optional
from urllib.parse import urlparse

logger = logging.getLogger('uvicorn.error')


# ─────────────────────────────────────────────
# Data contract
# ─────────────────────────────────────────────
@dataclass
class RawDocument:
    """
    Normalised output from every loader.
    text      : extracted text (may be empty if OCR needed)
    metadata  : source, page, type, language hint, etc.
    needs_ocr : True when the page is an image-only scan
    """
    text: str
    metadata: dict = field(default_factory=dict)
    needs_ocr: bool = False

    def is_empty(self) -> bool:
        return not self.text or not self.text.strip()


# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────
def _clean_text(text: str) -> str:
    """
    Normalize whitespace, fix common extraction artefacts,
    and keep Arabic + Latin unicode intact.
    """
    if not text:
        return ""
    # Normalize unicode (NFC keeps Arabic composites intact)
    text = unicodedata.normalize("NFC", text)
    # Remove null bytes and control chars (except newlines/tabs)
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", " ", text)
    # Collapse multiple blank lines → single blank line
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Collapse horizontal whitespace
    text = re.sub(r"[ \t]{2,}", " ", text)
    return text.strip()


def _detect_language_hint(text: str) -> str:
    """
    Fast heuristic — not a full detector.
    Returns 'arabic', 'english', or 'mixed'.
    """
    if not text:
        return "unknown"
    arabic = len(re.findall(r"[\u0600-\u06FF]", text))
    latin  = len(re.findall(r"[A-Za-z]",        text))
    total  = arabic + latin
    if total == 0:
        return "unknown"
    ratio = arabic / total
    if ratio > 0.65:
        return "arabic"
    if ratio < 0.25:
        return "english"
    return "mixed"


def _is_image_only_page(text: str) -> bool:
    """
    A page with fewer than 20 real characters after extraction
    is almost certainly a scanned image page.
    """
    return len(text.strip()) < 20


# ─────────────────────────────────────────────
# PDF Loader  (PyMuPDF — fitz)
# ─────────────────────────────────────────────
class PDFLoader:
    """
    Page-by-page PDF loader.
    Falls back to OCR flag when a page yields no text.
    Handles encrypted / corrupted files without crashing.
    
    NEW: Extracts embedded images from pages for OCR processing.
    """

    def load(self, file_path: str | Path) -> List[RawDocument]:
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF not installed. Run: pip install pymupdf")

        path = Path(file_path)
        docs: List[RawDocument] = []

        try:
            pdf = fitz.open(str(path))
        except Exception as exc:
            logger.error(f"[PDFLoader] Cannot open {path.name}: {exc}")
            return docs

        if pdf.is_encrypted:
            # Try empty password (many "protected" PDFs have no real password)
            if not pdf.authenticate(""):
                logger.warning(f"[PDFLoader] Skipping encrypted file: {path.name}")
                pdf.close()
                return docs

        for page_num in range(len(pdf)):
            try:
                page = pdf[page_num]
                raw_text = page.get_text("text")
                text = _clean_text(raw_text)
                needs_ocr = _is_image_only_page(text)

                # Extract images list for metadata (useful for OCR later)
                image_list = page.get_images(full=False)
                image_count = len(image_list)

                # Main page document
                docs.append(RawDocument(
                    text=text,
                    needs_ocr=needs_ocr,
                    metadata={
                        "source":       path.name,
                        "source_path":  str(path),
                        "source_type":  "pdf",
                        "page":         page_num + 1,
                        "total_pages":  len(pdf),
                        "has_images":   image_count > 0,
                        "language":     _detect_language_hint(text),
                    }
                ))
                
                # NEW: Extract embedded images for OCR
                # Even if page has text, images might contain additional info
                if image_count > 0 and not needs_ocr:
                    for img_idx, img_info in enumerate(image_list):
                        try:
                            xref = img_info[0]
                            base_image = pdf.extract_image(xref)
                            
                            if base_image:
                                # Create a separate document for this image
                                docs.append(RawDocument(
                                    text="",  # Will be filled by OCR
                                    needs_ocr=True,
                                    metadata={
                                        "source":       path.name,
                                        "source_path":  str(path),
                                        "source_type":  "pdf_embedded_image",
                                        "page":         page_num + 1,
                                        "total_pages":  len(pdf),
                                        "image_index":  img_idx + 1,
                                        "image_xref":   xref,
                                        "language":     "unknown",
                                    }
                                ))
                        except Exception as img_exc:
                            logger.debug(f"[PDFLoader] Could not extract image {img_idx} from page {page_num + 1}: {img_exc}")
                            
            except Exception as exc:
                logger.warning(f"[PDFLoader] Error on page {page_num + 1} of {path.name}: {exc}")

        pdf.close()
        logger.info(f"[PDFLoader] {path.name}: {len(docs)} documents, "
                    f"{sum(1 for d in docs if d.needs_ocr)} need OCR")
        return docs


# ─────────────────────────────────────────────
# Word Loader  (.docx)
# ─────────────────────────────────────────────
class WordLoader:
    """
    Extracts text from .docx files paragraph-by-paragraph.
    Preserves table cell content as structured text.
    NEW: Extracts embedded images for OCR processing.
    """

    def load(self, file_path: str | Path) -> List[RawDocument]:
        try:
            from docx import Document
            from PIL import Image
        except ImportError:
            raise ImportError("python-docx and Pillow not installed. Run: pip install python-docx Pillow")

        path = Path(file_path)
        parts: List[str] = []
        docs: List[RawDocument] = []

        try:
            doc = Document(str(path))
        except Exception as exc:
            logger.error(f"[WordLoader] Cannot open {path.name}: {exc}")
            return []

        # Paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())

        # Tables → convert to readable text
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        full_text = _clean_text("\n".join(parts))
        
        # Main document text
        if full_text:
            docs.append(RawDocument(
                text=full_text,
                metadata={
                    "source":      path.name,
                    "source_path": str(path),
                    "source_type": "docx",
                    "page":        1,
                    "language":    _detect_language_hint(full_text),
                }
            ))
        
        # NEW: Extract embedded images
        try:
            import io
            image_count = 0
            
            # Images are stored in document.part.rels
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_count += 1
                        image_data = rel.target_part.blob
                        
                        # Verify it's a valid image
                        img = Image.open(io.BytesIO(image_data))
                        
                        # Create a document for this image
                        docs.append(RawDocument(
                            text="",  # Will be filled by OCR
                            needs_ocr=True,
                            metadata={
                                "source":       path.name,
                                "source_path":  str(path),
                                "source_type":  "docx_embedded_image",
                                "page":         1,
                                "image_index":  image_count,
                                "image_data":   image_data,  # Store for OCR
                                "language":     "unknown",
                            }
                        ))
                    except Exception as img_exc:
                        logger.debug(f"[WordLoader] Could not extract image {image_count}: {img_exc}")
            
            if image_count > 0:
                logger.info(f"[WordLoader] {path.name}: Found {image_count} embedded images")
                
        except Exception as exc:
            logger.debug(f"[WordLoader] Image extraction failed (non-fatal): {exc}")

        if not docs:
            logger.warning(f"[WordLoader] No content extracted from {path.name}")
            
        return docs


# ─────────────────────────────────────────────
# PowerPoint Loader  (.pptx)
# ─────────────────────────────────────────────
class PowerPointLoader:
    """
    Extracts text from each slide as a separate document.
    Includes slide notes when available.
    NEW: Extracts embedded images for OCR processing.
    
    Extracts:
    - Text from shapes and placeholders
    - Embedded picture shapes
    - Images from grouped shapes
    - Background images from slides
    """

    def load(self, file_path: str | Path) -> List[RawDocument]:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from PIL import Image
        except ImportError:
            raise ImportError("python-pptx and Pillow not installed. Run: pip install python-pptx Pillow")

        path = Path(file_path)
        docs: List[RawDocument] = []

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            logger.error(f"[PowerPointLoader] Cannot open {path.name}: {exc}")
            return []

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: List[str] = []
            slide_image_count = 0

            # Extract from all shapes (including nested groups)
            def extract_from_shapes(shapes, depth=0):
                nonlocal slide_image_count
                
                for shape in shapes:
                    # Extract text from text-containing shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                    
                    # Extract picture shapes
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            # Check if shape has image data
                            if not hasattr(shape, 'image') or not shape.image:
                                logger.debug(f"[PowerPointLoader] Picture shape on slide {slide_num} has no image data")
                                continue
                                
                            image_data = shape.image.blob
                            
                            # Verify it's a valid image
                            import io
                            img = Image.open(io.BytesIO(image_data))
                            
                            slide_image_count += 1
                            
                            # Create a document for this image
                            docs.append(RawDocument(
                                text="",  # Will be filled by OCR
                                needs_ocr=True,
                                metadata={
                                    "source":       path.name,
                                    "source_path":  str(path),
                                    "source_type":  "pptx_embedded_image",
                                    "page":         slide_num,
                                    "total_pages":  len(prs.slides),
                                    "image_index":  slide_image_count,
                                    "image_data":   image_data,
                                    "language":     "unknown",
                                }
                            ))
                            logger.debug(f"[PowerPointLoader] Extracted image {slide_image_count} from slide {slide_num}")
                        except AttributeError as attr_exc:
                            logger.debug(f"[PowerPointLoader] Picture shape on slide {slide_num} missing image attribute: {attr_exc}")
                        except Exception as img_exc:
                            logger.debug(f"[PowerPointLoader] Could not extract image from slide {slide_num}: {img_exc}")
                    
                    # Recursively extract from grouped shapes
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        if hasattr(shape, 'shapes'):
                            logger.debug(f"[PowerPointLoader] Processing group with {len(shape.shapes)} shapes on slide {slide_num}")
                            extract_from_shapes(shape.shapes, depth + 1)
            
            # Extract from slide shapes
            extract_from_shapes(slide.shapes)
            
            # NEW: Try to extract background images
            try:
                # Check if slide has a background fill with an image
                if hasattr(slide, 'background') and hasattr(slide.background, 'fill'):
                    fill = slide.background.fill
                    if hasattr(fill, 'type') and fill.type == 6:  # 6 = picture fill
                        try:
                            # Access the image through the fill
                            if hasattr(fill, '_element'):
                                # This is a workaround to access background images
                                # Background images are harder to extract in python-pptx
                                logger.debug(f"[PowerPointLoader] Slide {slide_num} has background image (extraction not fully supported)")
                        except Exception as bg_exc:
                            logger.debug(f"[PowerPointLoader] Could not extract background image from slide {slide_num}: {bg_exc}")
            except Exception as exc:
                logger.debug(f"[PowerPointLoader] Background check failed for slide {slide_num}: {exc}")

            # Include presenter notes
            if slide.has_notes_slide:
                try:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Notes] {notes}")
                except Exception as notes_exc:
                    logger.debug(f"[PowerPointLoader] Could not extract notes from slide {slide_num}: {notes_exc}")

            text = _clean_text("\n".join(parts))
            
            # Add slide text document (even if empty, to maintain slide structure)
            docs.append(RawDocument(
                text=text,
                metadata={
                    "source":       path.name,
                    "source_path":  str(path),
                    "source_type":  "pptx",
                    "page":         slide_num,
                    "total_pages":  len(prs.slides),
                    "has_images":   slide_image_count > 0,
                    "language":     _detect_language_hint(text),
                }
            ))

        total_images = sum(1 for d in docs if d.needs_ocr)
        logger.info(f"[PowerPointLoader] {path.name}: {len(docs)} documents extracted, "
                   f"{total_images} images need OCR")
        return docs


# ─────────────────────────────────────────────
# Image Loader  (flags for OCR — no direct extraction)
# ─────────────────────────────────────────────
class ImageLoader:
    """
    Marks image files as needing OCR.
    Actual text extraction happens in ocr.py.
    """
    SUPPORTED = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp"}

    def load(self, file_path: str | Path) -> List[RawDocument]:
        path = Path(file_path)
        if path.suffix.lower() not in self.SUPPORTED:
            logger.warning(f"[ImageLoader] Unsupported image format: {path.suffix}")
            return []

        return [RawDocument(
            text="",
            needs_ocr=True,
            metadata={
                "source":      path.name,
                "source_path": str(path),
                "source_type": "image",
                "page":        1,
                "language":    "unknown",
            }
        )]


# ─────────────────────────────────────────────
# Text Loader  (.txt)
# ─────────────────────────────────────────────
class TextLoader:
    """
    Loads plain text files with encoding fallbacks.
    """

    ENCODINGS = ("utf-8", "utf-8-sig", "cp1252", "latin-1")

    def load(self, file_path: str | Path) -> List[RawDocument]:
        path = Path(file_path)

        for encoding in self.ENCODINGS:
            try:
                text = path.read_text(encoding=encoding, errors="strict")
                cleaned = _clean_text(text)

                if not cleaned:
                    logger.warning(f"[TextLoader] Empty text content: {path.name}")
                    return []

                return [RawDocument(
                    text=cleaned,
                    metadata={
                        "source":      path.name,
                        "source_path": str(path),
                        "source_type": "txt",
                        "page":        1,
                        "encoding":    encoding,
                        "language":    _detect_language_hint(cleaned),
                    }
                )]
            except UnicodeDecodeError:
                continue
            except Exception as exc:
                logger.error(f"[TextLoader] Cannot read {path.name}: {exc}")
                return []

        logger.error(f"[TextLoader] Could not decode file with supported encodings: {path.name}")
        return []


# ─────────────────────────────────────────────
# URL Loader
# ─────────────────────────────────────────────
class URLLoader:
    """
    Fetches a URL and extracts clean text.
    - Strips navigation, ads, scripts via BeautifulSoup
    - Respects timeout to avoid hanging
    - Handles redirects, encoding issues, and non-HTML responses
    """

    def __init__(self, timeout: int = 15):
        self.timeout = timeout

    def load(self, url: str) -> List[RawDocument]:
        try:
            import requests
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError(
                "requests and beautifulsoup4 not installed. "
                "Run: pip install requests beautifulsoup4"
            )

        # Validate URL
        parsed = urlparse(url)
        if parsed.scheme not in ("http", "https"):
            logger.error(f"[URLLoader] Invalid URL scheme: {url}")
            return []

        try:
            response = requests.get(
                url,
                timeout=self.timeout,
                headers={"User-Agent": "Mozilla/5.0 (RAG-Ingestion-Bot)"},
                allow_redirects=True,
            )
            response.raise_for_status()

            # Only process HTML
            content_type = response.headers.get("Content-Type", "")
            if "text/html" not in content_type:
                logger.warning(f"[URLLoader] Non-HTML content ({content_type}): {url}")
                return []

            # Detect encoding
            response.encoding = response.apparent_encoding or "utf-8"
            soup = BeautifulSoup(response.text, "html.parser")

            # Remove noise elements
            for tag in soup(["script", "style", "nav", "footer",
                              "header", "aside", "form", "iframe"]):
                tag.decompose()

            # Extract main content (prefer article/main over body)
            main = (soup.find("article") or soup.find("main") or soup.body)
            raw_text = main.get_text(separator="\n") if main else soup.get_text("\n")
            text = _clean_text(raw_text)

            if not text:
                logger.warning(f"[URLLoader] No text content extracted from: {url}")
                return []

            return [RawDocument(
                text=text,
                metadata={
                    "source":      parsed.netloc,
                    "source_path": url,
                    "source_type": "url",
                    "page":        1,
                    "language":    _detect_language_hint(text),
                    "status_code": response.status_code,
                }
            )]

        except Exception as exc:
            logger.error(f"[URLLoader] Failed to load {url}: {exc}")
            return []


# ─────────────────────────────────────────────
# Master Loader — dispatches by file type
# ─────────────────────────────────────────────
class DocumentLoader:
    """
    Single entry point for all document types.
    Auto-detects format and routes to the correct loader.
    """

    LOADERS = {
        ".txt":  TextLoader,
        ".pdf":  PDFLoader,
        ".docx": WordLoader,
        ".doc":  WordLoader,
        ".pptx": PowerPointLoader,
        ".ppt":  PowerPointLoader,
        ".png":  ImageLoader,
        ".jpg":  ImageLoader,
        ".jpeg": ImageLoader,
        ".tiff": ImageLoader,
        ".tif":  ImageLoader,
        ".bmp":  ImageLoader,
        ".webp": ImageLoader,
    }

    def load_file(self, file_path: str | Path) -> List[RawDocument]:
        path = Path(file_path)
        if not path.exists():
            logger.error(f"[DocumentLoader] File not found: {path}")
            return []

        ext = path.suffix.lower()
        loader_cls = self.LOADERS.get(ext)

        if loader_cls is None:
            logger.warning(f"[DocumentLoader] Unsupported file type: {ext}")
            return []

        loader = loader_cls()
        return loader.load(path)

    def load_url(self, url: str, timeout: int = 15) -> List[RawDocument]:
        return URLLoader(timeout=timeout).load(url)

    def load_directory(self, directory: str | Path,
                       recursive: bool = True) -> List[RawDocument]:
        """Load all supported files from a directory."""
        dir_path = Path(directory)
        if not dir_path.is_dir():
            logger.error(f"[DocumentLoader] Not a directory: {dir_path}")
            return []

        pattern = "**/*" if recursive else "*"
        all_docs: List[RawDocument] = []

        for file_path in dir_path.glob(pattern):
            if file_path.suffix.lower() in self.LOADERS:
                docs = self.load_file(file_path)
                all_docs.extend(docs)

        logger.info(f"[DocumentLoader] Directory scan: {len(all_docs)} documents loaded")
        return all_docs